"""
modules/file_engine.py – Gestione file universale per Cipher
                         Supporta: Excel, CSV, PDF, Word, PowerPoint, immagini, TXT, codice, AutoCAD (DXF)
                         Canali input: Telegram (uploads/) e Pi locale
"""

import os
import json
import shutil
from pathlib import Path
from typing import Optional
from datetime import datetime

from rich.console import Console
from config import Config

console = Console()
UPLOADS_DIR = Config.BASE_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)


class FileEngine:
    def __init__(self, llm_silent_fn=None) -> None:
        """
        llm_silent_fn: funzione Brain._call_llm_silent per descrizioni intelligenti
        """
        self._llm = llm_silent_fn

    # ── Router principale ─────────────────────────────────────────────────

    def process(self, path: str, instruction: str = "") -> str:
        """
        Punto di ingresso principale.
        Legge il file, lo analizza e risponde in base all'istruzione.
        """
        fpath = self._resolve_path(path)
        if not fpath.exists():
            return f"File non trovato: {path}"
        ext = fpath.suffix.lower()

        try:
            if ext in (".xlsx", ".xls"):
                return self._handle_excel(fpath, instruction)
            elif ext == ".csv":
                return self._handle_csv(fpath, instruction)
            elif ext == ".pdf":
                return self._handle_pdf(fpath, instruction)
            elif ext in (".docx", ".doc"):
                return self._handle_word(fpath, instruction)
            elif ext in (".pptx", ".ppt"):
                return self._handle_pptx(fpath, instruction)
            elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
                return self._handle_image(fpath, instruction)
            elif ext in (".dxf", ".dwg"):
                return self._handle_autocad(fpath, instruction)
            elif ext in (".txt", ".md", ".rst"):
                return self._handle_text(fpath, instruction)
            elif ext in (".py", ".js", ".ts", ".c", ".cpp", ".java", ".sh",
                         ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg",
                         ".html", ".css", ".sql", ".rs", ".go", ".al"):
                return self._handle_code(fpath, instruction)
            else:
                return self._handle_generic(fpath, instruction)
        except ImportError as e:
            return f"Libreria mancante per gestire {ext}: {e}\nInstalla con: pip install {self._suggest_package(ext)} --break-system-packages"
        except Exception as e:
            return f"Errore durante la lettura di {fpath.name}: {e}"

    # ── Excel ─────────────────────────────────────────────────────────────

    def _handle_excel(self, fpath: Path, instruction: str) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(fpath, data_only=True)
        result = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append([str(c) if c is not None else "" for c in row])
            if rows:
                result.append(f"### Foglio: {sheet_name}\n" + self._rows_to_text(rows))
        content = "\n\n".join(result)
        return self._apply_instruction(fpath, content, instruction, "excel")

    def _rows_to_text(self, rows: list) -> str:
        if not rows:
            return ""
        header = " | ".join(rows[0])
        sep    = "-" * len(header)
        data   = "\n".join(" | ".join(r) for r in rows[1:])
        return f"{header}\n{sep}\n{data}"

    # ── CSV ───────────────────────────────────────────────────────────────

    def _handle_csv(self, fpath: Path, instruction: str) -> str:
        import csv
        rows = []
        with fpath.open(encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        content = self._rows_to_text(rows[:50])  # max 50 righe
        if len(rows) > 50:
            content += f"\n... (e altre {len(rows)-50} righe)"
        return self._apply_instruction(fpath, content, instruction, "csv")

    # ── PDF ───────────────────────────────────────────────────────────────

    def _handle_pdf(self, fpath: Path, instruction: str) -> str:
        import fitz  # pymupdf
        doc  = fitz.open(str(fpath))
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        # Limita a 4000 caratteri per non sovraccaricare l'LLM
        if len(text) > 4000:
            text = text[:4000] + f"\n... (troncato, {len(text)} caratteri totali)"
        return self._apply_instruction(fpath, text, instruction, "pdf")

    # ── Word ──────────────────────────────────────────────────────────────

    def _handle_word(self, fpath: Path, instruction: str) -> str:
        from docx import Document
        doc  = Document(str(fpath))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if len(text) > 4000:
            text = text[:4000] + f"\n... (troncato)"
        return self._apply_instruction(fpath, text, instruction, "word")

    # ── PowerPoint ────────────────────────────────────────────────────────

    def _handle_pptx(self, fpath: Path, instruction: str) -> str:
        from pptx import Presentation
        prs   = Presentation(str(fpath))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        content = "\n".join(lines)
        if len(content) > 4000:
            content = content[:4000] + "\n... (troncato)"
        return self._apply_instruction(fpath, content, instruction, "powerpoint")

    # ── Immagini ──────────────────────────────────────────────────────────

    def _handle_image(self, fpath: Path, instruction: str) -> str:
        from PIL import Image
        img  = Image.open(fpath)
        info = f"Immagine: {fpath.name}\nDimensioni: {img.width}x{img.height}px\nFormato: {img.format}\nModalità: {img.mode}"

        if self._llm:
            prompt = (
                f"Descrivi brevemente questa immagine basandoti su questi metadati:\n{info}\n"
                f"Istruzione dell'utente: {instruction or 'descrivi il contenuto'}\n"
                f"Rispondi in italiano, 2-3 frasi."
            )
            try:
                return self._llm(prompt)
            except Exception:
                pass

        return info

    # ── AutoCAD ───────────────────────────────────────────────────────────

    def _handle_autocad(self, fpath: Path, instruction: str) -> str:
        if fpath.suffix.lower() == ".dwg":
            return "File DWG: formato binario proprietario. Converti in DXF per l'analisi."

        import ezdxf
        doc      = ezdxf.readfile(str(fpath))
        msp      = doc.modelspace()
        entities = {}
        for entity in msp:
            t = entity.dxftype()
            entities[t] = entities.get(t, 0) + 1

        layers = [layer.dxf.name for layer in doc.layers]
        summary = (
            f"File DXF: {fpath.name}\n"
            f"Layer: {', '.join(layers[:20])}\n"
            f"Entità: " + ", ".join(f"{k}={v}" for k, v in entities.items())
        )
        return self._apply_instruction(fpath, summary, instruction, "autocad")

    # ── Testo ─────────────────────────────────────────────────────────────

    def _handle_text(self, fpath: Path, instruction: str) -> str:
        content = fpath.read_text(encoding="utf-8", errors="replace")
        if len(content) > 4000:
            content = content[:4000] + f"\n... (troncato, {len(content)} caratteri totali)"
        return self._apply_instruction(fpath, content, instruction, "testo")

    # ── Codice ────────────────────────────────────────────────────────────

    def _handle_code(self, fpath: Path, instruction: str) -> str:
        content = fpath.read_text(encoding="utf-8", errors="replace")
        lines   = content.split("\n")
        if len(lines) > 200:
            content = "\n".join(lines[:200]) + f"\n... (troncato, {len(lines)} righe totali)"
        lang = fpath.suffix.lstrip(".")
        return self._apply_instruction(fpath, content, instruction, f"codice {lang}")

    # ── Generico ──────────────────────────────────────────────────────────

    def _handle_generic(self, fpath: Path, instruction: str) -> str:
        size = fpath.stat().st_size
        return (
            f"File: {fpath.name}\n"
            f"Tipo: {fpath.suffix}\n"
            f"Dimensione: {size} bytes\n"
            f"Formato non supportato per l'analisi diretta."
        )

    # ── Applica istruzione via LLM ────────────────────────────────────────

    def _apply_instruction(self, fpath: Path, content: str, instruction: str, file_type: str) -> str:
        """
        Se c'è un'istruzione specifica, chiede all'LLM di elaborare il contenuto.
        Altrimenti restituisce un riassunto.
        """
        if not instruction:
            instruction = "Fornisci un riassunto conciso del contenuto."

        if not self._llm:
            return f"Contenuto di {fpath.name}:\n\n{content}"

        prompt = (
            f"Sei Cipher. Hai letto un file {file_type}: '{fpath.name}'.\n"
            f"Contenuto:\n{content}\n\n"
            f"Istruzione: {instruction}\n\n"
            f"Rispondi in italiano, in modo diretto e conciso."
        )
        try:
            return self._llm(prompt)
        except Exception as e:
            return f"Contenuto letto ma errore LLM: {e}\n\n{content[:500]}"

    # ── Estrazione eventi calendario da Excel/CSV ─────────────────────────

    def extract_calendar_events(self, path: str) -> list[dict]:
        """
        Estrae eventi calendario da un file Excel o CSV.
        Parsing diretto senza LLM: legge data, fascia mattina, fascia pomeriggio.
        Ritorna lista di dict compatibili con calendar_create.
        """
        fpath = self._resolve_path(path)
        if not fpath.exists():
            return []

        events = []
        ext = fpath.suffix.lower()

        try:
            if ext in (".xlsx", ".xls"):
                import openpyxl
                from datetime import datetime as dt
                wb = openpyxl.load_workbook(fpath, data_only=True)
                ws = wb.active
                for row in ws.iter_rows(values_only=True):
                    date_val   = row[0]
                    mattina    = row[1]
                    pomeriggio = row[2]
                    # Salta righe senza data valida
                    if not date_val or not isinstance(date_val, dt):
                        continue
                    date_str = date_val.strftime("%Y-%m-%d")
                    tutor    = str(row[4]).strip() if row[4] else ""
                    location = str(row[5]).strip() if row[5] else ""
                    if mattina and str(mattina).strip():
                        parts = str(mattina).strip().split("-")
                        if len(parts) == 2:
                            events.append({
                                "title":       "Stage ITS - Mattina",
                                "start":       f"{date_str} {parts[0].strip()}",
                                "end":         f"{date_str} {parts[1].strip()}",
                                "description": f"Tutor: {tutor}",
                                "location":    location,
                            })
                    if pomeriggio and str(pomeriggio).strip():
                        parts = str(pomeriggio).strip().split("-")
                        if len(parts) == 2:
                            events.append({
                                "title":       "Stage ITS - Pomeriggio",
                                "start":       f"{date_str} {parts[0].strip()}",
                                "end":         f"{date_str} {parts[1].strip()}",
                                "description": f"Tutor: {tutor}",
                                "location":    location,
                            })
            elif ext == ".csv":
                import csv
                with fpath.open(encoding="utf-8", errors="replace") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        if len(row) < 3:
                            continue
                        date_str   = row[0].strip()
                        mattina    = row[1].strip()
                        pomeriggio = row[2].strip()
                        if not date_str:
                            continue
                        tutor    = row[4].strip() if len(row) > 4 else ""
                        location = row[5].strip() if len(row) > 5 else ""
                        if mattina:
                            parts = mattina.split("-")
                            if len(parts) == 2:
                                events.append({
                                    "title":       "Stage ITS - Mattina",
                                    "start":       f"{date_str} {parts[0].strip()}",
                                    "end":         f"{date_str} {parts[1].strip()}",
                                    "description": f"Tutor: {tutor}",
                                    "location":    location,
                                })
                        if pomeriggio:
                            parts = pomeriggio.split("-")
                            if len(parts) == 2:
                                events.append({
                                    "title":       "Stage ITS - Pomeriggio",
                                    "start":       f"{date_str} {parts[0].strip()}",
                                    "end":         f"{date_str} {parts[1].strip()}",
                                    "description": f"Tutor: {tutor}",
                                    "location":    location,
                                })
        except Exception as e:
            console.print(f"[red]extract_calendar_events error: {e}[/red]")
            return []

        return events

    # ── Modifica file ─────────────────────────────────────────────────────

    def modify_file(self, path: str, instruction: str) -> str:
        """Modifica un file in base all'istruzione."""
        fpath = self._resolve_path(path)
        if not fpath.exists():
            return f"File non trovato: {path}"

        ext = fpath.suffix.lower()

        if ext in (".txt", ".md", ".py", ".js", ".json", ".csv"):
            content = fpath.read_text(encoding="utf-8", errors="replace")
            if not self._llm:
                return "LLM non disponibile per la modifica."

            prompt = (
                f"Sei Cipher. Devi modificare questo file.\n"
                f"Contenuto attuale:\n{content}\n\n"
                f"Istruzione: {instruction}\n\n"
                f"Restituisci SOLO il contenuto modificato, nient'altro."
            )
            try:
                new_content = self._llm(prompt)
                fpath.write_text(new_content, encoding="utf-8")
                return f"File {fpath.name} modificato."
            except Exception as e:
                return f"Errore modifica: {e}"
        else:
            return f"Modifica diretta non supportata per {ext}. Posso solo leggere e analizzare questo formato."

    # ── Elimina file ──────────────────────────────────────────────────────

    def delete_file(self, path: str) -> str:
        """Elimina un file dalla cartella uploads o dal Pi."""
        fpath = self._resolve_path(path)
        if not fpath.exists():
            return f"File non trovato: {path}"
        try:
            fpath.unlink()
            return f"File {fpath.name} eliminato."
        except Exception as e:
            return f"Errore eliminazione: {e}"

    # ── Lista file uploads ────────────────────────────────────────────────

    def list_uploads(self) -> str:
        """Lista i file nella cartella uploads."""
        files = list(UPLOADS_DIR.iterdir())
        if not files:
            return "Nessun file nella cartella uploads."
        lines = []
        for f in sorted(files):
            size = f.stat().st_size
            lines.append(f"- {f.name} ({size} bytes)")
        return "File disponibili:\n" + "\n".join(lines)

    # ── Utility ───────────────────────────────────────────────────────────

    def _resolve_path(self, path: str) -> Path:
        """
        Risolve il path del file.
        Se è relativo, cerca prima in uploads/, poi nel filesystem del Pi.
        """
        p = Path(path)
        if p.is_absolute():
            return p
        # Cerca in uploads/
        upload_path = UPLOADS_DIR / path
        if upload_path.exists():
            return upload_path
        # Cerca relativo alla home di cipher
        cipher_path = Config.BASE_DIR / path
        if cipher_path.exists():
            return cipher_path
        # Path assoluto come ultimo tentativo
        return UPLOADS_DIR / path

    def _suggest_package(self, ext: str) -> str:
        mapping = {
            ".pdf":  "pymupdf",
            ".docx": "python-docx",
            ".doc":  "python-docx",
            ".pptx": "python-pptx",
            ".ppt":  "python-pptx",
            ".dxf":  "ezdxf",
            ".dwg":  "ezdxf",
        }
        return mapping.get(ext, "unknown")
